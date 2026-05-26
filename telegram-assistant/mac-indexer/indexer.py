#!/usr/bin/env python3
"""
indexer.py — Индексирует файлы Mac и rsync-ит индекс на VPS.

Запускается автоматически LaunchAgent каждые 2 часа.
Результат: ~/file-indexer/file_index.json → server:/home/parser/bots/assistant/data/
"""
import json
import os
import socket
import subprocess
import sys
from datetime import datetime
from pathlib import Path

# ── Настройки ─────────────────────────────────────────────────

SCAN_DIRS = [
    "~/Desktop",
    "~/Documents",
    "~/Downloads",
    "~/Library/Mobile Documents/com~apple~CloudDocs",  # iCloud Drive
]

# Папки, которые всегда включаем явно (даже если вложены в Desktop)
EXTRA_DIRS: list[str] = []

# Максимум символов из содержимого файла
SNIPPET_MAX = 600

# Расширения, из которых извлекаем текст
TEXT_EXTENSIONS = {".txt", ".md", ".py", ".js", ".ts", ".html", ".css", ".json",
                   ".yaml", ".yml", ".sh", ".csv", ".xml", ".ini", ".cfg", ".toml",
                   ".log", ".rst", ".tex"}
RICH_EXTENSIONS = {".pptx", ".docx", ".pdf", ".xlsx"}

# Папки, которые пропускаем
SKIP_DIRS = {
    ".git", ".svn", "node_modules", "__pycache__", ".cache",
    "Library", ".Trash", "Applications", "venv", ".venv",
}

# Куда сохраняем индекс
OUTPUT_DIR = Path("~/file-indexer").expanduser()
OUTPUT_FILE = OUTPUT_DIR / "file_index.json"
LOG_FILE = OUTPUT_DIR / "indexer.log"

# Сервер — SSH-алиас из ~/.ssh/config
REMOTE = "server:/home/parser/bots/assistant/data/file_index.json"


# ── Извлечение текста ──────────────────────────────────────────

def _extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    try:
        if ext in TEXT_EXTENSIONS:
            text = path.read_text(encoding="utf-8", errors="ignore")
            return text[:SNIPPET_MAX]

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                    if len("".join(parts)) >= SNIPPET_MAX:
                        break
            return " ".join(parts)[:SNIPPET_MAX]

        if ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            text = " ".join(p.text for p in doc.paragraphs if p.text.strip())
            return text[:SNIPPET_MAX]

        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                pages_text = []
                for page in pdf.pages[:5]:
                    t = page.extract_text()
                    if t:
                        pages_text.append(t)
                    if len("".join(pages_text)) >= SNIPPET_MAX:
                        break
            return " ".join(pages_text)[:SNIPPET_MAX]

        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets[:3]:
                for row in ws.iter_rows(max_row=20, values_only=True):
                    row_text = " ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        parts.append(row_text)
                    if len("".join(parts)) >= SNIPPET_MAX:
                        break
            return " ".join(parts)[:SNIPPET_MAX]

    except Exception:
        pass
    return ""


# ── Сканирование ───────────────────────────────────────────────

def scan_dir(base: Path) -> list[dict]:
    results = []
    if not base.exists():
        return results

    # mdfind (Spotlight) обходит ограничения macOS Privacy — rglob блокируется для LaunchAgent.
    # Ищем только файлы поддерживаемых форматов — wildcard '*' в mdfind не работает.
    all_exts = (
        list(TEXT_EXTENSIONS - {".log"})  # логи пропускаем — мусор
        + [".pptx", ".docx", ".pdf", ".xlsx"]
    )
    predicate = " || ".join(f"kMDItemFSName == '*{ext}'" for ext in all_exts)
    proc = subprocess.run(
        ["mdfind", "-onlyin", str(base), predicate],
        capture_output=True, text=True, timeout=120
    )
    entries = [Path(p.strip()) for p in proc.stdout.strip().split("\n") if p.strip()]

    for p in entries:
        # Пропускаем скрытые файлы и исключённые папки
        if any(part.startswith(".") for part in p.parts[-3:]):
            continue
        if any(skip in p.parts for skip in SKIP_DIRS):
            continue
        if not p.is_file():
            continue

        ext = p.suffix.lower()
        try:
            stat = p.stat()
            size_mb = round(stat.st_size / 1024 / 1024, 2)
            modified = datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds")
        except OSError:
            continue

        snippet = ""
        if ext in TEXT_EXTENSIONS or ext in RICH_EXTENSIONS:
            snippet = _extract_text(p)

        results.append({
            "name": p.name,
            "path": str(p),
            "ext": ext,
            "size_mb": size_mb,
            "modified": modified,
            "snippet": snippet,
        })
    return results


# ── Основной поток ─────────────────────────────────────────────

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # Пропускаем если индексировали менее 2 часов назад
    if OUTPUT_FILE.exists():
        try:
            last_str = json.loads(OUTPUT_FILE.read_text(encoding="utf-8")).get("indexed_at", "")
            if last_str:
                from datetime import timezone
                last_dt = datetime.fromisoformat(last_str)
                age_hours = (datetime.now() - last_dt).total_seconds() / 3600
                if age_hours < 2:
                    print(f"Индекс свежий ({age_hours:.1f}ч назад), пропускаем.")
                    return
        except Exception:
            pass

    log = open(LOG_FILE, "a", encoding="utf-8")

    def _log(msg: str):
        ts = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        line = f"[{ts}] {msg}"
        print(line)
        log.write(line + "\n")
        log.flush()

    _log("=== Индексация началась ===")

    dirs = [Path(d).expanduser() for d in SCAN_DIRS]
    dirs += [Path(d).expanduser() for d in EXTRA_DIRS]
    # убираем дубли
    seen = set()
    dirs = [d for d in dirs if not (str(d) in seen or seen.add(str(d)))]

    all_files = []
    for d in dirs:
        _log(f"Сканирую: {d}")
        files = scan_dir(d)
        _log(f"  → {len(files)} файлов")
        all_files.extend(files)

    # Убираем дубли по пути (EXTRA_DIRS могут пересекаться с SCAN_DIRS)
    seen_paths: set[str] = set()
    unique_files = []
    for f in all_files:
        if f["path"] not in seen_paths:
            seen_paths.add(f["path"])
            unique_files.append(f)

    index = {
        "indexed_at": datetime.now().isoformat(timespec="seconds"),
        "mac_hostname": socket.gethostname(),
        "total_files": len(unique_files),
        "files": unique_files,
    }

    OUTPUT_FILE.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    _log(f"Индекс сохранён: {OUTPUT_FILE} ({len(unique_files)} файлов, {OUTPUT_FILE.stat().st_size // 1024} KB)")

    # Синхронизация на сервер
    _log(f"rsync → {REMOTE}")
    result = subprocess.run(
        ["rsync", "-az", str(OUTPUT_FILE), REMOTE],
        capture_output=True, text=True,
    )
    if result.returncode == 0:
        _log("rsync: OK")
    else:
        _log(f"rsync ОШИБКА: {result.stderr.strip()}")
        sys.exit(1)

    _log("=== Готово ===\n")
    log.close()


if __name__ == "__main__":
    main()
